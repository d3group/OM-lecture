import sys
from pptx import Presentation
import os

def extract_pptx_text(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return

    prs = Presentation(pptx_path)
    
    print(f"--- Content of {os.path.basename(pptx_path)} ---")
    
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1}:")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"- {shape.text}")

if __name__ == "__main__":
    pptx_path = "/Users/duypham/Documents/OM-lecture/MRP Logic.pptx"
    extract_pptx_text(pptx_path)
